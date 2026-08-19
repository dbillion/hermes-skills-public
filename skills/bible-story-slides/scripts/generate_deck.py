#!/usr/bin/env python3
"""
Generate a black-and-white, kid-friendly (ages 6-9) Bible story slide deck.

Usage:
    python3 generate_deck.py --story-json story.json --out deck.pptx
"""
import argparse
import json
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def add_text(slide, text, left, top, width, height, size, bold=False,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = BLACK
    return box


def add_illustration_placeholder(slide, left, top, width, height, label):
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = BLACK
    frame.line.width = Pt(3)
    tf = frame.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ illustration: {label} ]"
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.name = FONT
    run.font.color.rgb = BLACK


def title_slide(prs, title, big_idea):
    slide = blank_slide(prs)
    add_text(slide, title, Inches(1), Inches(2.3), Inches(11.3), Inches(1.6),
              size=60, bold=True)
    add_text(slide, big_idea, Inches(1.5), Inches(4.1), Inches(10.3), Inches(1.2),
              size=28)
    return slide


def beat_slide(prs, index, total, beat_text):
    slide = blank_slide(prs)
    add_illustration_placeholder(
        slide, Inches(2.67), Inches(0.6), Inches(8), Inches(3.6),
        f"scene {index} of {total}"
    )
    add_text(slide, beat_text, Inches(1.2), Inches(4.5), Inches(10.9), Inches(2.2),
              size=44, bold=True)
    add_text(slide, f"{index} / {total}", Inches(12.2), Inches(6.9),
              Inches(1.0), Inches(0.5), size=16, align=PP_ALIGN.RIGHT)
    return slide


def list_slide(prs, heading, items):
    slide = blank_slide(prs)
    add_text(slide, heading, Inches(1), Inches(0.7), Inches(11.3), Inches(1.1),
              size=40, bold=True)
    box = slide.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(10.3), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"\u2022 {item}"
        run.font.size = Pt(32)
        run.font.name = FONT
        run.font.color.rgb = BLACK
        p.space_after = Pt(18)
    return slide


def build_deck(story, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, story["title"], story["big_idea"])

    beats = story["beats"]
    for i, beat in enumerate(beats, start=1):
        beat_slide(prs, i, len(beats), beat)

    list_slide(prs, "What can we learn?", story["takeaways"])
    list_slide(prs, "Let's talk about it", story["discussion_questions"])

    prs.save(out_path)
    print(f"Wrote {out_path} ({2 + len(beats) + 0} content slides + title)")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--story-json", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    with open(args.story_json, "r", encoding="utf-8") as f:
        story = json.load(f)

    required = ["title", "big_idea", "beats", "takeaways", "discussion_questions"]
    missing = [k for k in required if k not in story or not story[k]]
    if missing:
        sys.exit(f"story.json missing required fields: {missing}")

    build_deck(story, args.out)


if __name__ == "__main__":
    main()
